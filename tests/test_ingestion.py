import pytest
from unittest.mock import patch, MagicMock
from core.ingestion import ingest_file, delete_document_chunks, _load_pptx


def test_unsupported_file_type_raises(tmp_path):
    fake_file = tmp_path / "notes.xyz"
    fake_file.write_text("content")
    with pytest.raises(ValueError, match="Unsupported file type"):
        ingest_file(str(fake_file), "course-1", "notes.xyz", "xyz", "notes", "doc-1")


def test_delete_document_chunks_noop_when_vectorstore_missing():
    delete_document_chunks("nonexistent-course-xyz", "nonexistent-doc-xyz")


def _make_ingest_mocks():
    """Return (mock_chroma_class, mock_vectorstore) with metadata-supporting chunk."""
    mock_vs = MagicMock()
    mock_chroma = MagicMock(return_value=mock_vs)
    return mock_chroma, mock_vs


@patch("core.ingestion.chromadb.PersistentClient")
@patch("core.ingestion.Chroma")
@patch("core.ingestion.HuggingFaceEmbeddings")
@patch("core.ingestion.RecursiveCharacterTextSplitter")
@patch("core.ingestion.TextLoader")
def test_ingest_txt_returns_chunk_count(mock_loader, mock_splitter, mock_embeddings, mock_chroma, mock_client, tmp_path):
    fake_file = tmp_path / "notes.txt"
    fake_file.write_text("Some study notes.")

    fake_doc = MagicMock()
    fake_doc.metadata = {}
    mock_loader.return_value.load.return_value = [fake_doc]

    chunks = [MagicMock(metadata={}) for _ in range(3)]
    mock_splitter.return_value.split_documents.return_value = chunks

    mock_vs = MagicMock()
    mock_chroma.return_value = mock_vs

    count = ingest_file(str(fake_file), "course-1", "notes.txt", "txt", "notes", "doc-1")

    assert count == 3
    mock_vs.add_documents.assert_called_once()


@patch("core.ingestion.chromadb.PersistentClient")
@patch("core.ingestion.Chroma")
@patch("core.ingestion.HuggingFaceEmbeddings")
@patch("core.ingestion.RecursiveCharacterTextSplitter")
@patch("core.ingestion.PyPDFLoader")
def test_ingest_pdf_uses_pdf_loader(mock_loader, mock_splitter, mock_embeddings, mock_chroma, mock_client, tmp_path):
    fake_file = tmp_path / "lecture.pdf"
    fake_file.write_bytes(b"%PDF fake")

    fake_doc = MagicMock()
    fake_doc.metadata = {}
    mock_loader.return_value.load.return_value = [fake_doc]
    mock_splitter.return_value.split_documents.return_value = [MagicMock(metadata={})]
    mock_chroma.return_value = MagicMock()

    ingest_file(str(fake_file), "c1", "lecture.pdf", "pdf", "slides", "d1")

    mock_loader.assert_called_once_with(str(fake_file))


@patch("core.ingestion.chromadb.PersistentClient")
@patch("core.ingestion.Chroma")
@patch("core.ingestion.HuggingFaceEmbeddings")
@patch("core.ingestion.RecursiveCharacterTextSplitter")
@patch("core.ingestion.TextLoader")
def test_chunk_metadata_is_set(mock_loader, mock_splitter, mock_embeddings, mock_chroma, mock_client, tmp_path):
    fake_file = tmp_path / "notes.txt"
    fake_file.write_text("content")

    fake_doc = MagicMock()
    fake_doc.metadata = {}
    mock_loader.return_value.load.return_value = [fake_doc]

    chunk = MagicMock()
    chunk.metadata = {}
    mock_splitter.return_value.split_documents.return_value = [chunk]
    mock_chroma.return_value = MagicMock()

    ingest_file(str(fake_file), "cid", "notes.txt", "txt", "notes", "did")

    assert chunk.metadata["source"] == "notes.txt"
    assert chunk.metadata["category"] == "notes"
    assert chunk.metadata["course_id"] == "cid"
    assert chunk.metadata["document_id"] == "did"


def test_load_pptx_extracts_slide_content(tmp_path):
    """_load_pptx returns one Document per slide that has text."""
    from pptx import Presentation
    prs = Presentation()
    layout = prs.slide_layouts[1]
    for i in range(3):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"Title {i + 1}"
        slide.placeholders[1].text = f"Body text on slide {i + 1}."
    pptx_path = tmp_path / "deck.pptx"
    prs.save(str(pptx_path))

    docs = _load_pptx(str(pptx_path))

    assert len(docs) == 3
    for i, doc in enumerate(docs):
        assert doc.metadata["page"] == i + 1
        assert f"slide {i + 1}" in doc.page_content.lower()


def test_load_pptx_skips_empty_slides(tmp_path):
    """Slides with no text are not returned."""
    from pptx import Presentation
    prs = Presentation()
    empty_layout = prs.slide_layouts[6]  # blank layout — no text placeholders
    prs.slides.add_slide(empty_layout)
    pptx_path = tmp_path / "empty.pptx"
    prs.save(str(pptx_path))

    docs = _load_pptx(str(pptx_path))
    assert docs == []


@patch("core.ingestion.chromadb.PersistentClient")
@patch("core.ingestion.Chroma")
@patch("core.ingestion.HuggingFaceEmbeddings")
@patch("core.ingestion.RecursiveCharacterTextSplitter")
@patch("core.ingestion._load_pptx")
def test_ingest_pptx_uses_load_pptx(mock_load_pptx, mock_splitter, mock_embeddings, mock_chroma, mock_client, tmp_path):
    fake_file = tmp_path / "slides.pptx"
    fake_file.write_bytes(b"fake pptx")

    fake_doc = MagicMock()
    fake_doc.metadata = {}
    mock_load_pptx.return_value = [fake_doc]
    mock_splitter.return_value.split_documents.return_value = [MagicMock(metadata={})]
    mock_chroma.return_value = MagicMock()

    count = ingest_file(str(fake_file), "c1", "slides.pptx", "pptx", "slides", "d1")

    mock_load_pptx.assert_called_once_with(str(fake_file))
    assert count == 1
