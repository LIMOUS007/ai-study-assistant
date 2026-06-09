import shutil
from pathlib import Path
from langchain_community.document_loaders import PyPDFLoader, TextLoader, Docx2txtLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
import httpx
import chromadb
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import Chroma
import requests
from youtube_transcript_api import YouTubeTranscriptApi
from langchain_core.documents import Document
from pptx import Presentation


def _get_vectorstore(course_id: str) -> Chroma:
    """Create a Chroma instance using PersistentClient to avoid the shared-system-client bug."""
    vectorstore_path = Path("vectorstore") / course_id
    vectorstore_path.mkdir(parents=True, exist_ok=True)
    client = chromadb.PersistentClient(path=str(vectorstore_path))
    # embeddings = OpenAIEmbeddings(http_client=httpx.Client(verify=False))
    # embeddings = GoogleGenerativeAIEmbeddings(model="models/gemini-embedding-001")
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    return Chroma(client=client, embedding_function=embeddings)


def _load_pptx(file_path: str) -> list[Document]:
    prs = Presentation(file_path)
    docs = []
    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        if texts:
            docs.append(Document(
                page_content="\n".join(texts),
                metadata={"page": slide_num},
            ))
    return docs


def ingest_file(file_path, course_id, filename, file_type, document_category, document_id) -> int:
    suffix = Path(file_path).suffix.lower()
    if suffix == ".txt":
        loader = TextLoader(file_path, encoding="utf-8")
        docs = loader.load()
    elif suffix == ".pdf":
        loader = PyPDFLoader(file_path)
        docs = loader.load()
    elif suffix == ".docx":
        loader = Docx2txtLoader(file_path)
        docs = loader.load()
    elif suffix == ".pptx":
        docs = _load_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")

    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    chunks = splitter.split_documents(docs)

    for chunk in chunks:
        chunk.metadata.update({
            "source": filename,
            "category": document_category,
            "course_id": course_id,
            "document_id": document_id,
        })

    vector_store = _get_vectorstore(course_id)
    vector_store.add_documents(chunks)
    return len(chunks)


def delete_vectorstore(course_id: str):
    """Delete the entire per-course ChromaDB vector store."""
    import gc
    vectorstore_path = Path("vectorstore") / course_id
    if not vectorstore_path.exists():
        return
    try:
        client = chromadb.PersistentClient(path=str(vectorstore_path))
        for col in client.list_collections():
            client.delete_collection(col.name)
        del client
        gc.collect()
    except Exception:
        pass
    shutil.rmtree(vectorstore_path, ignore_errors=True)


def delete_document_chunks(course_id: str, document_id: str):
    """Remove all Chroma chunks that belong to a specific document."""
    vectorstore_path = Path("vectorstore") / course_id
    if not vectorstore_path.exists():
        return
    vector_store = _get_vectorstore(course_id)
    result = vector_store._collection.get(where={"document_id": {"$eq": document_id}})
    if result["ids"]:
        vector_store.delete(result["ids"])


def _video_id_from_url(url: str) -> str:
    import re
    match = re.search(r"(?:v=|youtu\.be/)([A-Za-z0-9_-]{11})", url)
    return match.group(1) if match else url


def _fetch_video_title(url: str, session: requests.Session) -> str:
    """Fetch the video title via YouTube's oEmbed API — no API key or pytube needed."""
    try:
        resp = session.get(
            f"https://www.youtube.com/oembed?url={url}&format=json",
            timeout=5,
        )
        if resp.status_code == 200:
            return resp.json().get("title", "")
    except Exception:
        pass
    return ""


def ingest_youtube(url: str, course_id: str, document_id: str) -> tuple[str, int]:
    """
    Load YouTube transcript -> split -> embed -> store in Chroma.
    Returns (video_title, chunk_count).

    Raises TranscriptsDisabled or NoTranscriptFound from youtube_transcript_api
    if captions are unavailable -- let the caller handle those for the UI error message.
    """
    video_id = _video_id_from_url(url)

    # SSL verification disabled -- required on networks with SSL inspection.
    session = requests.Session()
    session.verify = False

    video_title = _fetch_video_title(url, session)
    title = video_title if video_title else f"YouTube: {video_id}"

    api = YouTubeTranscriptApi(http_client=session)
    transcript = api.fetch(video_id)
    full_text = " ".join(snippet.text for snippet in transcript)

    docs = [Document(page_content=full_text)]
    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    chunks = splitter.split_documents(docs)

    for chunk in chunks:
        chunk.metadata.update({
            "source": title,
            "category": "youtube",
            "course_id": course_id,
            "document_id": document_id,
            "page": 0,
        })

    vector_store = _get_vectorstore(course_id)
    vector_store.add_documents(chunks)
    return title, len(chunks)
